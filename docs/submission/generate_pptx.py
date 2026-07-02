from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "春招平行宇宙\nCareer Multiverse"
    subtitle.text = "让 AI 替你跑完 1000 次春招\n\n可实验的人才市场沙盘\ntljcpa\n2026年7月"

    # Slide 2: The Problem
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "应届生求职的真问题：三个看不见"
    content = slide.placeholders[1]
    content.text = "每年1158万应届生，人均投递60+简历，70%没有回音。\n\n" \
                   "1. 看不见 HR 的真实评价：淘汰理由不透明\n" \
                   "2. 看不见 公司的隐性门槛：实际门槛远高于JD\n" \
                   "3. 看不见 决策的后悔成本：反事实无法验证\n\n" \
                   "现状：人才市场至今只能凭经验决策，缺一个“先实验、再决策”的地方。"

    # Slide 3: The Solution
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "解决方案：把整个春招“虚拟跑一遍”"
    content = slide.placeholders[1]
    content.text = "1. 上传真实简历 + GitHub + Blog\n" \
                   "2. LLM 一次性多任务评估（五维画像评分+学校学历档加成）\n" \
                   "3. 数字分身进入 3D 沙盘宇宙\n" \
                   "   - 49家公司独立 HR Agent + 200名竞争者\n" \
                   "   - 蒙特卡洛扩展跑完 1000 次平行春招\n" \
                   "4. 输出预演结果：\n" \
                   "   - 1000次平行宇宙的KPI统计（offer率、去向分布）\n" \
                   "   - 反事实预演（如果改XX，结局如何变化）\n" \
                   "   - AI教练个性化可执行建议"

    # Slide 4: Innovations
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "三个核心反转"
    content = slide.placeholders[1]
    content.text = "1. 结构反转：把一次性的投递赌博，变成可重复的1000次完整春招预演\n" \
                   "2. 载体反转：从静态PDF简历，变成可交互面试的“数字分身”\n" \
                   "3. 决策反转：不是含糊地“建议你怎么做”，而是直接展示不同选择在1000次模拟中的“统计结果”\n\n" \
                   "竞品都在局部优化，我们在重新定义市场结构。"

    # Slide 4b: 宝妈绝境案例（真实推演，诚实呈现"没翻盘 = 赛道错配"）
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "绝境案例：一位宝妈的真实推演（真数据·未美化）"
    content = slide.placeholders[1]
    content.text = "画像：28岁 / 文科本科 / 2年全职带娃空窗 / 目标项目管理\n" \
                   "在约300家公司市场里，每档各真跑8次（共24次真实LLM沙盘）：\n\n" \
                   "① 残酷现实（原始）：平均投52份 → 0面试 · 0 offer（8/8全军覆没）\n" \
                   "② 只改包装（带娃重构为多线程+主攻创业公司+薪资-10%）：仍 0 offer（8/8）\n" \
                   "③ 真补硬实力（PMP拿证+两个真实项目，含金量25→50）：本赛道仍 0 offer（8/8）\n\n" \
                   "沙盘诚实的结论：她的瓶颈不是简历不够漂亮，是赛道错配——\n" \
                   "该市场约2/3岗位以技术岗打头、平均门槛≈80、超半数要求硕士，她根本够不着。\n" \
                   "真正的“后悔药”是趁早换一条能上桌的赛道，而非优化措辞。\n" \
                   "（数字全部真跑得到，结果就是“没翻盘”，如实呈现——诚实是产品立身之本）"

    # Slide 5: Tech Stack
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "技术架构"
    content = slide.placeholders[1]
    content.text = "前端：Vue3 + Three.js (3D沙盘) + D3.js (可视化)\n" \
                   "后端：FastAPI + Python 3.10\n" \
                   "核心引擎：自研轻量 Multi-Agent 编排（不用LangGraph）\n" \
                   " - LLM tier路由，适配多种国内大模型\n" \
                   " - Candidate/CompanyHR/Interviewer Agent\n" \
                   "部署方案：\n" \
                   " - Azure VM + Caddy + Cloudflare CDN"

    # Slide 6: Business Model
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "商业模式：主战场是高校"
    content = slide.placeholders[1]
    content.text = "主营收：高校就业指导中心（to G SaaS订阅）\n" \
                   " - 购买群体级人才洞察，将就业指导从“经验驱动”升级为“数据驱动”\n\n" \
                   "获客层：C端应届生（免费+数据飞轮）\n" \
                   " - 口碑传播渗透校内，反推学校采购\n\n" \
                   "第二曲线：B端企业与招聘平台\n" \
                   " - 企业反向品牌报告（你公司在应届生眼里的真实画像）\n" \
                   " - 与招聘平台集成"

    # Slide 6b: Digital Twin Calibration Loop (VISION / ROADMAP - explicitly NOT implemented yet)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "数字孪生对齐引擎（愿景/路线图，尚未实现）"
    content = slide.placeholders[1]
    content.text = "【本页为接入智联真实数据后的规划，当前尚未实现】\n\n" \
                   "现状（冷启动）：\n" \
                   " - 沙盘里约300个HR Agent的偏好完全由LLM先验知识驱动\n" \
                   " - 尚未被任何一条真实招聘行为数据校准过\n\n" \
                   "规划（接入智联数据后）：\n" \
                   " - 接入平台真实行为数据（已读不回/发起沟通/发放offer）\n" \
                   " - 用强化学习持续校准HR Agent的偏好分布\n" \
                   " - 使沙盘从孤立模拟升级为招聘市场的“实时数字孪生”\n\n" \
                   "诚实边界：这一步没有回测数据支撑效果，讲的是方向，不是成果——" \
                   "也是我们希望通过本次决赛促成的合作方向。"

    # Slide 7: Compliance
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "数据合规（0风险）"
    content = slide.placeholders[1]
    content.text = "所有数据均为合成数据：\n" \
                   " - 49家虚构公司全部代号化（无真实名称）\n" \
                   " - 200名求职者均为化名生成的persona\n\n" \
                   "双层兜底拦截：\n" \
                   " - Prompt约束不准出现真名\n" \
                   " - 落盘前强字典扫描替换（47家真公司名拦截）\n\n" \
                   "零PII处理：\n" \
                   " - 用户上传简历仅驻留会话内存，不上云不留存"

    # Slide 8: Team and Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "团队与未来展望"
    content = slide.placeholders[1]
    content.text = "团队构成：\n" \
                   " - tljcpa (跨考非科班学生)\n" \
                   " - Claude Code (AI 辅助结对)\n" \
                   " - 这是AI赋能下，1个人+AI完成完整产品开发的真实样本\n\n" \
                   "未来诉求：\n" \
                   " - 期望加入智联生态合作\n" \
                   " - 寻求种子轮投资（100-200万）\n" \
                   " - 校园渠道推广与公开数据合作"

    output_path = '/root/智联AI比赛/docs/submission/春招平行宇宙-决赛答辩.pptx'
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"PPTX created successfully at {output_path}")

if __name__ == '__main__':
    create_presentation()
