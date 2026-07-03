"""
决赛答辩 PPT 生成(python-pptx)。
和 deepseek 文本堆版的区别:每页以真页面截图 / 设计图为主视觉，文字克制，深色统一。
素材：/tmp/ppt-shots/ 下 7 张 demo 真截图 + 5 张设计图；二维码在 docs/submission/。
输出：docs/submission/春招平行宇宙-决赛答辩.pptx
本文件生成后，wly 在 PowerPoint 打开过一遍、微调、导出 PDF 备用。

新增第 10c 页「数字孪生对齐引擎」：明确标注为接入智联真实数据后的规划/路线图，
当前尚未实现，不构成已交付能力——诚实边界写进了 bullet 原文，讲解时不可省略这层措辞。

新增第 12 页「宝妈绝境案例」：虚构 persona（合成典型画像，非真人，零 PII）真跑真数据
（scripts/mom_comeback_sim.py），诚实呈现"没翻盘=赛道错配"，严禁讲成翻盘/拿offer的励志故事——
沙盘的价值是提前告诉你残酷真相，不是保证人人能赢。
本轮跑的时候 DeepSeek 额度耗尽，临时切到 NVIDIA 免费层的 llama-3.1-8b-instruct（弱模型），
24 次全 0 offer 的结果可能部分反映弱模型筛选偏严苛，bullet 措辞已改为"沙盘推演揭示的方向"，
不写成板上钉钉的确定结论——讲解时这层限定不能省略。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SHOTS="/tmp/ppt-shots"
QR="/root/智联AI比赛/docs/submission/demo-二维码-multiverse.png"
OUT="/root/智联AI比赛/docs/submission/春招平行宇宙-决赛答辩.pptx"

BG=RGBColor(0x0A,0x0E,0x1A); PANEL=RGBColor(0x13,0x18,0x29)
CYAN=RGBColor(0x22,0xD3,0xEE); PURPLE=RGBColor(0xA8,0x55,0xF7)
GOLD=RGBColor(0xFB,0xBF,0x24); INK=RGBColor(0xE5,0xE9,0xF0); INK2=RGBColor(0x8B,0x93,0xA7)
FONT="微软雅黑"
SW,SH=Inches(13.333),Inches(7.5)

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
BLANK=prs.slide_layouts[6]

def bg_dark(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=BG

def full_img(img):
    s=prs.slides.add_slide(BLANK); bg_dark(s)
    s.shapes.add_picture(img,0,0,width=SW,height=SH)  # 设计图已是完整16:9视觉，铺满
    return s

def tbox(s,x,y,w,h,text,size,color,bold=True,align=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.color.rgb=color; f.name=FONT
    return tb

def shot_page(img,title,sub=""):
    s=prs.slides.add_slide(BLANK); bg_dark(s)
    tbox(s,Inches(0.6),Inches(0.35),Inches(12),Inches(0.8),title,26,CYAN)
    # 截图作为主视觉，居中留边（16:9 → 宽 10.6 高 5.96）
    iw=Inches(10.6); ih=Inches(5.96); ix=(SW-iw)//2; iy=Inches(1.25)
    s.shapes.add_picture(img,ix,iy,width=iw,height=ih)
    if sub: tbox(s,Inches(0.6),Inches(7.0),Inches(12),Inches(0.4),sub,13,INK2,bold=False)
    return s

def text_page(title,bullets,accent=CYAN):
    s=prs.slides.add_slide(BLANK); bg_dark(s)
    tbox(s,Inches(0.9),Inches(0.9),Inches(11.5),Inches(1.0),title,32,accent)
    # 卡片
    card=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.9),Inches(2.3),Inches(11.5),Inches(4.2))
    card.fill.solid(); card.fill.fore_color.rgb=PANEL; card.line.color.rgb=CYAN; card.line.width=Pt(1)
    card.shadow.inherit=False
    tf=card.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.5); tf.margin_top=Inches(0.4)
    tf.vertical_anchor=MSO_ANCHOR.TOP
    for i,b in enumerate(bullets):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(14)
        r=p.add_run(); r.text="·  "+b; f=r.font; f.size=Pt(18); f.color.rgb=INK; f.name=FONT; f.bold=False
    return s

def end_page(title,sub):
    s=prs.slides.add_slide(BLANK); bg_dark(s)
    tbox(s,Inches(1),Inches(2.2),Inches(11.3),Inches(1.2),title,34,INK,align=PP_ALIGN.CENTER)
    tbox(s,Inches(1),Inches(3.5),Inches(11.3),Inches(0.7),sub,20,CYAN,bold=False,align=PP_ALIGN.CENTER)
    # 二维码 + 说明
    qw=Inches(1.8); s.shapes.add_picture(QR,(SW-qw)//2,Inches(4.5),width=qw,height=qw)
    tbox(s,Inches(1),Inches(6.4),Inches(11.3),Inches(0.5),"扫码体验 multiverse.zdwktlj.top  ·  github.com/tljcpa/career-multiverse",13,INK2,bold=False,align=PP_ALIGN.CENTER)
    return s

# ===== 16 页（原 14 页 + 数字孪生愿景页 + 宝妈绝境案例页 + 身位差 bullet 已并入第 4 页）=====
full_img(f"{SHOTS}/design_cover.png")                                    # 1 封面
full_img(f"{SHOTS}/design_problem.png")                                  # 2 问题
shot_page(f"{SHOTS}/report.png","把整个春招反复推演，用统计结论替代单次赌博","offer 率 98.4% / 平均 5.3 offer / 中位薪资 53 万 / 第 10 周落定 / 去向分散在声剧-γ·鲸火-1·云枢等 8+ 家——基于数十次真实模拟的统计结论，引擎可放大到全量真实市场")  # 3
text_page("从预测未来，到实验未来",[                                       # 4
 "现有 AI 求职工具：当前数据 → 一个预测结果（简历改词 / 面试外挂 / 词袋匹配）",
 "我们：当前状态 → 构建虚拟市场 → 观察演化 → 反事实分析",
 "从推荐系统到模拟系统的跃迁——改变一个条件，看结局怎么变",
 "身位差：别人打磨更快的镰刀，我们造气候模拟器——工具改变效率，只有沙盘能改变战略",
])
full_img(f"{SHOTS}/design_threelayer.png")                               # 5 三层视角
shot_page(f"{SHOTS}/profile.png","LLM 五维画像评估，透明不黑盒","每维附评分理由（引用简历证据）+ 学校档 / 学历档双维度加成，综合分可展开")  # 6
shot_page(f"{SHOTS}/sandbox.png","3D 沙盘：约 300 家公司星团 · 你的数字分身","每家公司独立 HR Agent，可现场采访；点击公司节点触发投递")  # 7
text_page("反事实预演：从预测未来，到实验未来（核心卖点）",[  # 8
 "拖动三个滑块：简历质量 / 项目含金量 / 学校档，系统把市场重新推演一遍",
 "输出 offer 数、薪资、去向的变化——例：项目含金量 +15 → 平均 offer 从 5.3 升到 9.4、中位薪资从 53 万升到 62 万",
 "不止告诉你「简历 80 分」，更告诉你「改了会怎样」",
 "把职业规划从玄学，变成基于概率与因果的科学决策（现场 demo 可动态演示）",
])
full_img(f"{SHOTS}/design_arch.png")                                     # 9 架构
full_img(f"{SHOTS}/design_chain.png")                                    # 10 商业价值链
text_page("从 demo 到全量：我们的引擎 × 平台的数据",[                    # 10b 规模化 / 合作
 "现在：约 300 家虚构公司的沙盘——作用是证明引擎跑得通",
 "核心资产不是数据量，是「让人才市场可实验」的引擎与方法论",
 "接入平台级全量真实数据 + 算力 → 覆盖真实市场的人才市场数字孪生",
 "这正是与平台最自然的合作：我们出引擎和方法，平台出数据、算力、场景",
 "「缺数据」不是短板，是合作接口；全量下用分层架构（粗筛 + 精模拟）承接",
])
text_page("数字孪生对齐引擎（愿景 / 路线图，尚未实现）",[                 # 10c 诚实愿景页
 "【本页讲的是接入智联真实数据后的规划，不是已经做到的事】",
 "现状 · 冷启动：约 300 个 HR Agent 的偏好完全由 LLM 先验知识驱动，尚未被任何一条真实招聘数据校准过",
 "规划 · 接入数据后：用智联平台每日真实行为数据（已读不回 / 发起沟通 / 发放 offer）+ 强化学习持续校准 HR Agent",
 "愿景：沙盘从孤立的模拟，升级为招聘市场的「实时数字孪生」",
 "诚实边界：这一步没有回测数据支撑效果，讲的是方向而非成果——也是我们想促成与智联合作的原因",
],accent=GOLD)
text_page("人力资本市场的量化回测沙盘",[                          # 11 商业定位升维 + 先验证
 "金融投资 10 万都先跑一万次回测；择业这种人生最大投资、企业几百万招聘成本却仍靠肉身试错",
 "接入真实数据后，这里可以是招聘界的「彭博终端」：C 端职业生涯 GPS 卖确定性，"
 "B 端雇主品牌与薪酬动态定价沙盘（调薪前先跑一遍，看人才流失率变化，省下真实试错成本，想象空间最大）",
 "我们判断高校可能最先跑通（有就业率 KPI + 预算），但这是待验证的假设，不是算好的账",
 "早期该做的是找标杆校验证真实付费意愿、转起数据飞轮，而非纸面财务预测",
],accent=GOLD)
text_page("绝境案例：一个虚构 persona 的沙盘推演（合成画像 · 未美化）",[    # 12 诚实反例：赛道错配
 "画像：虚构 persona（合成典型画像，非真人，零PII）——28岁/双非本科文科/2年全职带娃空窗/目标项目管理，在约300家公司大市场真跑，每档各8次（共24次真实LLM沙盘）",
 "数据说明：本轮因 DeepSeek 额度耗尽，临时切到 NVIDIA 免费层 llama-3.1-8b-instruct（弱模型）完成，24次全0的结果可能部分反映弱模型筛选偏严苛，以下是这轮推演揭示的方向，非确定结论",
 "① 基线（原始简历）：平均投52份 → 0面试 · 0 offer（8/8）",
 "② 反事实·纯包装（带娃重构为多线程管理+主攻创业公司+薪资降10%）：仍 0 offer（8/8）——包装未能改变硬维度筛选结果",
 "③ 反事实·真补实力（PMP拿证+两个真实项目，含金量25→50）：本赛道仍 0 offer（8/8）",
 "推演方向：瓶颈可能不是简历措辞，而是赛道错配——该市场307家里263家以技术岗打头、57%要求硕士，双非文科转岗门槛明显更高",
 "沙盘不是用来证明「努力就能翻盘」的，而是提前给出一个诚实的方向：与其在同一条赛道死磕包装，不如尽早评估换赛道的可能性",
],accent=GOLD)
text_page("非科班单人 · 22 天 · AI 当队友",[                              # 13
 "创作者 tljcpa：机械原理跨考考研中，非计算机科班",
 "AI（Claude Code）做工程实现，核心创意 / 产品方向 / 架构决策由本人主导",
 "落地证据：真实录屏演示 + 开源代码（可本地部署）+ 线上可访问 demo",
],accent=PURPLE)
end_page("求职不该是一场只有一次的赌博","让人才市场——先实验，再决策；不要用你唯一真实的人生去试错")        # 14

prs.save(OUT)
print("生成完成：",OUT,"共",len(prs.slides._sldIdLst),"页")
